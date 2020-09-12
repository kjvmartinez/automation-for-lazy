from pptx import Presentation
import csv

def sendNamesToTemplate():
    # open template file
    prs = Presentation("/Users/kjvmartinez/Documents/Dev/Active_Projects/question-to-pptx/certificate_tmp.pptx")
    title_slide_layout = prs.slide_layouts[0]

    with open("participants.csv") as file:
        reader = csv.reader(file)
        for row in reader:
            slide = prs.slides.add_slide(title_slide_layout)
            print(str(reader.line_num) + " - " + str(row[0]))
            p = prs.slides[reader.line_num-1].shapes[0]
            p.text = row[0]
            
    prs.save('certificate_tmp.pptx')


sendNamesToTemplate()
print("Done.")