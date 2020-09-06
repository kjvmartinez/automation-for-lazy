from pptx import Presentation
import csv

def sendQsToPpt():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]

    with open("get_questions.csv") as file:
        reader = csv.reader(file)
        next(reader)
        for row in reader:
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            # print(repr(row))
            title.text = row[0]

    prs.save('foo.pptx')

sendQsToPpt()
print("Done.")